import asyncio
import os
import g4f
import re
from aiogram import Bot, Dispatcher, types
from aiogram.filters import Command
from aiogram.types import FSInputFile
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- ТОКЕН ---
API_TOKEN = '8584797614:AAGax0agzl31zVzpjWp19DeFYdehuUAMP2M'

bot = Bot(token=API_TOKEN)
dp = Dispatcher()

async def get_ai_content(topic, lang_code):
    prompts = {
        "en": "English",
        "uz": "O'zbek tili",
        "ru": "Русский язык"
    }
    target_lang = prompts.get(lang_code, prompts["ru"])
    
    # Максимально жесткий промпт, чтобы ИИ не халтурил
    full_prompt = (
        f"Create a HUGE educational presentation on: '{topic}'. Language: {target_lang}. "
        f"You MUST write 10 slides. "
        f"Structure: Slide Title; Slide Content (minimum 10-15 long sentences) || next slide... "
        f"Use '||' as a separator between EVERY slide. This is CRITICAL. "
        f"Each slide must have at least 1000 characters of text. And maximum 2400 characters"
        f"No greetings, no 'Slide 1' text. Just raw content."
    )
    
    try:
        response = await g4f.ChatCompletion.create_async(
            model=g4f.models.default,
            messages=[{"role": "user", "content": full_prompt}],
        )
        # Если ИИ забыл про разделители, пробуем найти их по смыслу (запасной вариант)
        if "||" not in response and ";" in response:
            response = response.replace("\n\n", "||")
            
        return response
    except Exception as e:
        return f"Error;{e}"

def create_final_pptx(topic, ai_text, filename):
    prs = Presentation('template.pptx') if os.path.exists('template.pptx') else Presentation()

    # Очистка текста от мусора
    clean_text = re.sub(r'(?i)(слайд|slide|slayd)\s*\d*\s*[:\-]*', '', ai_text).replace("**", "")
    
    # Разбиваем на слайды и убираем пустые куски
    slides_raw = [s.strip() for s in clean_text.split('||') if len(s.strip()) > 10]
    
    if not slides_raw: # Если разделителей нет, пробуем делить по двойному переносу строки
        slides_raw = [s.strip() for s in clean_text.split('\n\n') if len(s.strip()) > 10]

    for i, data in enumerate(slides_raw[:10]):
        if ';' in data:
            parts = data.split(';')
            title_t = parts[0].strip()
            body_t = ";".join(parts[1:]).strip() # Собираем остаток текста, если там были лишние ";"
        else:
            title_t = topic if i == 0 else f"Slide {i+1}"
            body_t = data

        layout_idx = 6 if len(prs.slide_layouts) > 6 else 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # --- ТИТУЛЬНЫЙ ЛИСТ ---
        if i == 0:
            f_size = Pt(36) if len(title_t) < 50 else Pt(26)
            t_box = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(2.5))
            tf = t_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_t
            p.alignment = PP_ALIGN.CENTER
            p.font.size = f_size
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            continue

        # --- ОСТАЛЬНЫЕ СЛАЙДЫ ---
        h_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(1))
        h_tf = h_box.text_frame
        h_tf.text = title_t
        h_p = h_tf.paragraphs[0]
        h_p.font.size = Pt(24)
        h_p.font.bold = True
        h_p.font.color.rgb = RGBColor(255, 255, 255)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.4))
        b_tf = body_box.text_frame
        b_tf.word_wrap = True
        b_tf.auto_size = 1 
        
        p = b_tf.add_paragraph()
        p.text = body_t
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

    prs.save(filename)

@dp.message()
async def handle_all(message: types.Message):
    raw_text = message.text
    if not raw_text or raw_text.startswith('/'): return

    lang = "ru"
    if "/en" in raw_text: lang = "en"; raw_text = raw_text.replace("/en", "")
    elif "/uz" in raw_text: lang = "uz"; raw_text = raw_text.replace("/uz", "")
    
    topic = raw_text.strip()
    # Чистим имя файла от запрещенных символов
    safe_name = re.sub(r'[\\/*?:"<>|]', "", topic[:20])
    file_path = f"{safe_name or 'Presentation'}.pptx"

    status = await message.answer(f"⏳ Собираю информацию на языке: {lang.upper()}...")
    
    content = await get_ai_content(topic, lang)
    
    try:
        create_final_pptx(topic, content, file_path)
        await message.answer_document(FSInputFile(file_path), caption=f"✅ Готово! Найдено листов: {content.count('||') + 1}")
        os.remove(file_path)
    except Exception as e:
        await message.answer(f"Ошибка при создании: {e}")
    finally:
        await status.delete()

async def main():
    await dp.start_polling(bot)

if __name__ == '__main__':
    asyncio.run(main())